from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = Path(__file__).resolve().parent / "mini-ctf_weekly_report.pptx"

NAVY = RGBColor(8, 14, 26)
PANEL = RGBColor(16, 27, 45)
PANEL2 = RGBColor(22, 38, 62)
BLUE = RGBColor(99, 169, 255)
PALE = RGBColor(224, 239, 255)
MUTED = RGBColor(168, 188, 211)
WHITE = RGBColor(248, 251, 255)
GREEN = RGBColor(114, 220, 170)
ORANGE = RGBColor(255, 183, 107)


def textbox(slide, text, x, y, w, h, size=18, color=WHITE, bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.04)
    frame.margin_right = Inches(0.04)
    frame.margin_top = Inches(0.02)
    p = frame.paragraphs[0]
    p.text = text
    p.font.name = "Malgun Gothic"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    if align:
        p.alignment = align
    return box


def panel(slide, x, y, w, h, fill=PANEL, line=PANEL2, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1)
    return shape


def bullet_box(slide, items, x, y, w, h, size=16, color=PALE, gap=8):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.08)
    for index, item in enumerate(items):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.text = f"• {item}"
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(gap)
    return box


def slide_base(prs, number, title, kicker):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = NAVY
    textbox(slide, f"MINI / CTF  ·  WEEKLY REPORT  ·  {number:02d}", 0.65, 0.35, 6.5, 0.25, 9, BLUE, True)
    textbox(slide, kicker.upper(), 0.65, 0.82, 11.5, 0.28, 10, MUTED, True)
    textbox(slide, title, 0.65, 1.12, 12.1, 0.65, 29, WHITE, True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(1.92), Inches(12.0), Inches(0.015))
    line.fill.solid(); line.fill.fore_color.rgb = BLUE; line.line.fill.background()
    textbox(slide, "Mini-CTF MVP  ·  PostgreSQL / Spring Boot / React", 0.65, 7.08, 8, 0.2, 8, MUTED)
    textbox(slide, str(number), 12.15, 7.08, 0.5, 0.2, 8, BLUE, True, PP_ALIGN.RIGHT)
    return slide


def add_stat(slide, value, label, x, y, w, color=BLUE):
    panel(slide, x, y, w, 1.02, PANEL, PANEL2, True)
    textbox(slide, value, x + 0.18, y + 0.12, w - 0.3, 0.38, 25, color, True)
    textbox(slide, label, x + 0.18, y + 0.58, w - 0.3, 0.24, 11, MUTED)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = NAVY
    textbox(slide, "MINI / CTF", 0.75, 0.72, 5.6, 0.5, 22, BLUE, True)
    textbox(slide, "보안 학습과 문제 풀이를 연결한\n온라인 CTF 플랫폼", 0.75, 1.55, 10.8, 1.45, 34, WHITE, True)
    textbox(slide, "주간 개발 보고  ·  MVP 구현 현황", 0.8, 3.3, 6, 0.35, 17, MUTED)
    panel(slide, 8.25, 1.15, 3.7, 3.7, PANEL, BLUE, True)
    textbox(slide, "{  }", 9.0, 1.78, 2.2, 0.8, 57, BLUE, True, PP_ALIGN.CENTER)
    textbox(slide, "learn by breaking\nthings safely", 8.8, 3.0, 2.6, 0.7, 16, PALE, True, PP_ALIGN.CENTER)
    textbox(slide, "2026.08  ·  Mini-CTF Organization", 0.8, 6.75, 6, 0.25, 10, MUTED)

    slide = slide_base(prs, 2, "프로젝트 목표와 현재 범위", "01  ·  Product overview")
    textbox(slide, "Mini-CTF는 문제 풀이만 제공하는 서비스가 아니라, 인증·학습·커뮤니티·운영을 하나로 묶은 CTF 플랫폼입니다.", 0.8, 2.25, 11.5, 0.5, 19, PALE)
    add_stat(slide, "3", "단계별 CTF 문제", 0.8, 3.25, 2.7, BLUE)
    add_stat(slide, "6", "Flyway DB migrations", 3.75, 3.25, 2.7, GREEN)
    add_stat(slide, "2", "OAuth providers live", 6.7, 3.25, 2.7, ORANGE)
    add_stat(slide, "1", "Admin operations console", 9.65, 3.25, 2.7, BLUE)
    bullet_box(slide, ["사용자가 직접 문제를 내려받고 풀이·제출", "관리자가 사용자·문제·로그를 운영", "풀이와 커뮤니티 활동을 통한 학습 경험 제공"], 0.9, 4.75, 11.4, 1.2, 16)

    slide = slide_base(prs, 3, "시스템 아키텍처", "02  ·  Technical foundation")
    for x, title, sub, color in [(0.9, "React + Vite", "Frontend / Theme / UX", BLUE), (4.85, "Spring Boot", "REST API / Security", GREEN), (8.8, "PostgreSQL", "JPA / Flyway V1~V6", ORANGE)]:
        panel(slide, x, 2.75, 3.0, 1.45, PANEL, color, True)
        textbox(slide, title, x + 0.18, 3.05, 2.65, 0.32, 20, color, True, PP_ALIGN.CENTER)
        textbox(slide, sub, x + 0.18, 3.55, 2.65, 0.25, 12, PALE, False, PP_ALIGN.CENTER)
    textbox(slide, "REST API", 3.95, 3.2, 0.75, 0.25, 11, MUTED, True, PP_ALIGN.CENTER)
    textbox(slide, "JPA / SQL", 7.93, 3.2, 0.75, 0.25, 11, MUTED, True, PP_ALIGN.CENTER)
    textbox(slide, "SSE 실시간 메시지", 4.85, 4.95, 3.0, 0.3, 14, BLUE, True, PP_ALIGN.CENTER)
    bullet_box(slide, ["Swagger/OpenAPI로 API 확인 가능", "PostgreSQL 16 연결 및 Flyway 스키마 검증", "실행 주소: Frontend 5173 / Backend 8080"], 1.0, 5.55, 11, 0.9, 15)

    slide = slide_base(prs, 4, "회원가입·로그인과 보안", "03  ·  Authentication")
    panel(slide, 0.85, 2.3, 5.7, 3.75, PANEL, PANEL2, True)
    textbox(slide, "구현 완료", 1.15, 2.62, 2.3, 0.32, 19, GREEN, True)
    bullet_box(slide, ["회원가입 / 일반 로그인", "Argon2 비밀번호 해시", "JWT 인증 및 사용자 세션", "Google OAuth 로그인", "GitHub OAuth 로그인", "OAuth 실패 메시지 처리"], 1.05, 3.1, 5.1, 2.5, 15)
    panel(slide, 6.9, 2.3, 5.5, 3.75, PANEL, PANEL2, True)
    textbox(slide, "운영 보완", 7.2, 2.62, 2.3, 0.32, 19, ORANGE, True)
    bullet_box(slide, ["이메일 인증", "비밀번호 찾기 / 재설정", "로그인 시도 잠금", "운영 도메인 OAuth Redirect URI", "HTTPS 및 Secure Cookie"], 7.1, 3.1, 4.9, 2.5, 15)

    slide = slide_base(prs, 5, "CTF 문제 풀이 코어", "04  ·  Challenge engine")
    headers = ["난이도", "문제", "분야", "학습 목표", "점수"]
    rows = [["Easy", "Signal in Plain Sight", "Crypto", "Base64 디코딩", "100"], ["Medium", "Proxy Afterimage", "Forensics", "Hex → Base64", "250"], ["Hard", "Orbit Gatekeeper", "Reversing", "변환식 역연산", "500"]]
    x0, y0, widths = 0.8, 2.45, [1.1, 3.35, 1.55, 3.1, 1.0]
    cursor = x0
    for head, width in zip(headers, widths):
        panel(slide, cursor, y0, width, 0.5, PANEL2, BLUE); textbox(slide, head, cursor + 0.1, y0 + 0.12, width - 0.2, 0.2, 12, PALE, True, PP_ALIGN.CENTER); cursor += width
    for ri, row in enumerate(rows):
        cursor = x0
        for ci, (value, width) in enumerate(zip(row, widths)):
            panel(slide, cursor, y0 + 0.55 + ri * 0.65, width, 0.6, PANEL, PANEL2)
            textbox(slide, value, cursor + 0.1, y0 + 0.73 + ri * 0.65, width - 0.2, 0.2, 12 if ci != 1 else 13, WHITE if ci == 1 else PALE, ci == 1, PP_ALIGN.CENTER)
            cursor += width
    bullet_box(slide, ["문제 상세 조회 및 Artifact 다운로드", "FLAG 제출·오답 기록·점수 반영", "동시 정답 제출 시 중복 점수 방지", "3번 검증기의 signed/unsigned byte 버그 수정 완료"], 1.0, 5.2, 11.0, 1.0, 15)

    slide = slide_base(prs, 6, "안티치트와 관리자 운영", "05  ·  Trust & safety")
    panel(slide, 0.85, 2.3, 5.65, 3.8, PANEL, PANEL2, True)
    textbox(slide, "행동 기반 보조 탐지", 1.15, 2.63, 3.8, 0.35, 19, BLUE, True)
    bullet_box(slide, ["문제 활동 없이 정답 제출", "문제 접근 직후의 빠른 정답", "오답 제출 및 빈도 제한", "문제 접근·Artifact 다운로드 이력", "Easy 20초 / Medium 60초 / Hard 120초"], 1.05, 3.15, 5.0, 2.5, 15)
    panel(slide, 6.8, 2.3, 5.65, 3.8, PANEL, PANEL2, True)
    textbox(slide, "Admin Console", 7.1, 2.63, 3.8, 0.35, 19, GREEN, True)
    bullet_box(slide, ["사용자 수정·정지·복구", "문제 CRUD 및 Artifact 관리", "제출·Anti-cheat 이벤트 조회", "로그 Redact / Hide", "관리자 감사 로그"], 7.0, 3.15, 5.0, 2.5, 15)
    textbox(slide, "AI 사용 여부를 단정하는 탐지기가 아니라, 검토 가능한 보안 신호를 제공하는 구조", 1.0, 6.45, 11.2, 0.3, 14, ORANGE, True, PP_ALIGN.CENTER)

    slide = slide_base(prs, 7, "프로필·커뮤니티·실시간 메시지", "06  ·  Social learning")
    cards = [(0.85, "Profile", ["닉네임·상태 메시지", "클릭형 프로필 사진 업로드"], BLUE), (4.45, "Community", ["게시글·댓글 CRUD", "작성자·관리자 권한 제어"], GREEN), (8.05, "Social", ["친구 요청·수락", "SSE 실시간 개인 메시지"], ORANGE)]
    for x, title, items, color in cards:
        panel(slide, x, 2.65, 3.25, 2.5, PANEL, color, True)
        textbox(slide, title, x + 0.2, 2.98, 2.8, 0.35, 20, color, True)
        bullet_box(slide, items, x + 0.15, 3.65, 2.9, 1.0, 15)
    textbox(slide, "문제 풀이 중심 서비스에서 사용자 간 학습·토론이 가능한 플랫폼으로 확장", 1.0, 5.8, 11.1, 0.45, 20, PALE, True, PP_ALIGN.CENTER)

    slide = slide_base(prs, 8, "프론트엔드 UX 개선", "07  ·  User experience")
    add_stat(slide, "Dark / Light", "테마 전환", 0.85, 2.45, 2.8, BLUE)
    add_stat(slide, "High contrast", "문제 제출 영역", 3.9, 2.45, 2.8, GREEN)
    add_stat(slide, "Responsive", "화면 대응", 6.95, 2.45, 2.8, ORANGE)
    add_stat(slide, "Branding", "Mini-CTF SVG 로고", 10.0, 2.45, 2.5, BLUE)
    bullet_box(slide, ["라이트 테마에서 글자·테두리 대비 강화", "작은 설명·랭킹·마이페이지 폰트 크기 개선", "Google·GitHub 버튼 및 로고 스타일 개선", "프로필 이미지 클릭 업로드 흐름 단순화", "오류·로딩·빈 상태 화면 추가 개선 예정"], 1.0, 4.35, 11, 1.45, 16)

    slide = slide_base(prs, 9, "검증 결과와 현재 상태", "08  ·  Verification")
    add_stat(slide, "PASS", "Backend compile / test", 0.85, 2.45, 2.8, GREEN)
    add_stat(slide, "PASS", "Frontend build / lint", 3.9, 2.45, 2.8, GREEN)
    add_stat(slide, "PASS", "PostgreSQL + Flyway", 6.95, 2.45, 2.8, GREEN)
    add_stat(slide, "PASS", "3 Artifact verifiers", 10.0, 2.45, 2.5, GREEN)
    bullet_box(slide, ["Backend API: http://localhost:8080", "Frontend: http://localhost:5173", "Swagger: /swagger-ui.html", "GitHub main 브랜치 최신 코드 반영", "실제 사용자 시나리오 기반 E2E 테스트는 추가 예정"], 1.0, 4.35, 11, 1.35, 16)

    slide = slide_base(prs, 10, "주차별 개발 진행 계획", "09  ·  Weekly roadmap")
    weeks = [
        (0.78, "1주차", "기반 구축", ["Spring Boot + React 구성", "PostgreSQL·Flyway 연결", "회원가입·일반 로그인"], BLUE),
        (3.88, "2주차", "CTF 핵심 기능", ["Easy/Medium/Hard 3문제", "FLAG 제출·점수·랭킹", "Admin·Anti-cheat 신호"], GREEN),
        (6.98, "3주차", "서비스 확장", ["Google/GitHub OAuth", "프로필·커뮤니티·댓글", "친구·실시간 메시지·테마"], ORANGE),
        (10.08, "4주차", "예정: MVP 안정화", ["E2E·보안 테스트", "HTTPS·OAuth 운영 설정", "백업·모니터링·배포"], BLUE),
    ]
    for x, week, title, items, color in weeks:
        panel(slide, x, 2.45, 2.55, 3.5, PANEL, color, True)
        textbox(slide, week, x + 0.18, 2.76, 2.15, 0.25, 13, color, True)
        textbox(slide, title, x + 0.18, 3.15, 2.15, 0.5, 17, WHITE, True)
        bullet_box(slide, items, x + 0.12, 4.05, 2.25, 1.45, 13, PALE, 6)
    textbox(slide, "현재까지 1~3주차 핵심 기능 구현 완료 · 4주차는 운영 가능한 MVP 전환 단계", 0.95, 6.35, 11.2, 0.35, 17, PALE, True, PP_ALIGN.CENTER)

    slide = slide_base(prs, 11, "정리", "10  ·  Closing")
    textbox(slide, "핵심 기능은 구현되었습니다.", 0.85, 2.4, 6.5, 0.55, 28, WHITE, True)
    textbox(slide, "다음 단계는 운영 안정성·보안·콘텐츠 품질을 높여\n실제 사용 가능한 MVP로 완성하는 것입니다.", 0.85, 3.25, 8.6, 0.9, 22, PALE)
    panel(slide, 9.5, 2.25, 2.45, 2.45, PANEL, BLUE, True)
    textbox(slide, "MINI\n/ CTF", 9.8, 2.82, 1.85, 1.0, 28, BLUE, True, PP_ALIGN.CENTER)
    textbox(slide, "learn by breaking things safely", 0.85, 6.28, 5.5, 0.25, 12, MUTED)

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
